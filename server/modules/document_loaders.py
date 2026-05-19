"""Multi-format document loaders for SU-GPT (Section 2).

Each loader returns a list of normalized records:
    {
        "text": str,
        "metadata": {
            "source": str,           # filename
            "document_type": str,    # pdf | pptx | docx | md | txt
            "page": int | None,      # PDFs
            "slide": int | None,     # PPTX
            "section": str | None,   # DOCX heading / MD heading
        }
    }
"""

from __future__ import annotations

import json
from pathlib import Path
from typing import Any, Callable

from langchain_community.document_loaders import PyPDFLoader


SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".pptx",
    ".docx",
    ".md",
    ".txt",
    ".html",
    ".htm",
    ".ipynb",
    ".json",
}


def _record(text: str, source: str, document_type: str, **extra: Any) -> dict:
    metadata = {
        "source": source,
        "document_type": document_type,
        "page": extra.get("page"),
        "slide": extra.get("slide"),
        "section": extra.get("section"),
    }
    return {"text": text, "metadata": metadata}


def load_pdf(file_path: str) -> list[dict]:
    source = Path(file_path).name
    loader = PyPDFLoader(file_path)
    docs = loader.load()
    records = []
    for doc in docs:
        # PyPDFLoader uses 0-indexed pages in metadata['page']; expose 1-indexed
        raw_page = doc.metadata.get("page")
        page = (raw_page + 1) if isinstance(raw_page, int) else None
        records.append(_record(doc.page_content, source, "pdf", page=page))
    return records


def load_pptx(file_path: str) -> list[dict]:
    from pptx import Presentation  # imported lazily so missing dep doesn't break PDFs

    source = Path(file_path).name
    prs = Presentation(file_path)
    records = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                line = "".join(run.text for run in paragraph.runs)
                if line.strip():
                    parts.append(line)
        text = "\n".join(parts).strip()
        if text:
            records.append(_record(text, source, "pptx", slide=slide_index))
    return records


def load_docx(file_path: str) -> list[dict]:
    from docx import Document  # imported lazily

    source = Path(file_path).name
    doc = Document(file_path)

    records: list[dict] = []
    current_section: str | None = None
    buffer: list[str] = []

    def flush() -> None:
        if not buffer:
            return
        text = "\n".join(buffer).strip()
        if text:
            records.append(_record(text, source, "docx", section=current_section))
        buffer.clear()

    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style_name = (paragraph.style.name or "").lower() if paragraph.style else ""
        if style_name.startswith("heading"):
            flush()
            current_section = text
        else:
            buffer.append(text)
    flush()

    if not records:
        # Fall back to a single record if no headings/paragraphs were captured
        full_text = "\n".join(p.text for p in doc.paragraphs).strip()
        if full_text:
            records.append(_record(full_text, source, "docx"))

    return records


def _read_text_file(file_path: str) -> str:
    path = Path(file_path)
    try:
        return path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return path.read_text(encoding="utf-8", errors="replace")


def load_markdown(file_path: str) -> list[dict]:
    source = Path(file_path).name
    raw = _read_text_file(file_path)

    records: list[dict] = []
    current_section: str | None = None
    buffer: list[str] = []

    def flush() -> None:
        if not buffer:
            return
        text = "\n".join(buffer).strip()
        if text:
            records.append(_record(text, source, "md", section=current_section))
        buffer.clear()

    for line in raw.splitlines():
        stripped = line.lstrip()
        if stripped.startswith("#"):
            heading = stripped.lstrip("#").strip()
            flush()
            current_section = heading or current_section
        else:
            buffer.append(line)
    flush()

    if not records and raw.strip():
        records.append(_record(raw, source, "md"))

    return records


def load_txt(file_path: str) -> list[dict]:
    source = Path(file_path).name
    text = _read_text_file(file_path).strip()
    if not text:
        return []
    return [_record(text, source, "txt")]


def load_html(file_path: str) -> list[dict]:
    from bs4 import BeautifulSoup  # imported lazily

    source = Path(file_path).name
    raw = _read_text_file(file_path)
    soup = BeautifulSoup(raw, "html.parser")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()

    records: list[dict] = []
    current_section: str | None = None
    buffer: list[str] = []

    def flush() -> None:
        if not buffer:
            return
        text = "\n".join(buffer).strip()
        if text:
            records.append(_record(text, source, "html", section=current_section))
        buffer.clear()

    for node in soup.find_all(["h1", "h2", "h3", "h4", "p", "li", "pre", "code"]):
        text = node.get_text(" ", strip=True)
        if not text:
            continue
        if node.name in {"h1", "h2", "h3", "h4"}:
            flush()
            current_section = text
        else:
            buffer.append(text)
    flush()

    if not records:
        text = soup.get_text("\n", strip=True)
        if text:
            records.append(_record(text, source, "html"))
    return records


def load_ipynb(file_path: str) -> list[dict]:
    source = Path(file_path).name
    notebook = json.loads(_read_text_file(file_path))
    records: list[dict] = []
    current_section: str | None = None
    buffer: list[str] = []

    def cell_source(cell: dict[str, Any]) -> str:
        raw = cell.get("source", "")
        if isinstance(raw, list):
            return "".join(str(part) for part in raw).strip()
        return str(raw).strip()

    def flush() -> None:
        if not buffer:
            return
        text = "\n\n".join(buffer).strip()
        if text:
            records.append(_record(text, source, "ipynb", section=current_section))
        buffer.clear()

    for cell in notebook.get("cells", []):
        if not isinstance(cell, dict):
            continue
        text = cell_source(cell)
        if not text:
            continue
        cell_type = cell.get("cell_type", "cell")
        first_line = text.splitlines()[0].strip()
        if cell_type == "markdown" and first_line.startswith("#"):
            flush()
            current_section = first_line.lstrip("#").strip() or current_section
            remaining = "\n".join(text.splitlines()[1:]).strip()
            if remaining:
                buffer.append(remaining)
        else:
            prefix = "Code example" if cell_type == "code" else "Notebook note"
            buffer.append(f"{prefix}:\n{text}")
    flush()

    if not records:
        text = "\n\n".join(
            cell_source(cell)
            for cell in notebook.get("cells", [])
            if isinstance(cell, dict)
        ).strip()
        if text:
            records.append(_record(text, source, "ipynb"))
    return records


def load_json(file_path: str) -> list[dict]:
    source = Path(file_path).name
    data = json.loads(_read_text_file(file_path))
    records: list[dict] = []

    if isinstance(data, list):
        group_size = 25
        for start in range(0, len(data), group_size):
            group = data[start : start + group_size]
            text = json.dumps(group, ensure_ascii=False, indent=2)
            records.append(
                _record(
                    text,
                    source,
                    "json",
                    section=f"items {start + 1}-{start + len(group)}",
                )
            )
    elif isinstance(data, dict):
        for key, value in data.items():
            text = json.dumps(value, ensure_ascii=False, indent=2)
            records.append(_record(text, source, "json", section=str(key)))
    else:
        records.append(_record(json.dumps(data, ensure_ascii=False, indent=2), source, "json"))
    return records


_LOADERS: dict[str, Callable[[str], list[dict]]] = {
    ".pdf": load_pdf,
    ".pptx": load_pptx,
    ".docx": load_docx,
    ".md": load_markdown,
    ".txt": load_txt,
    ".html": load_html,
    ".htm": load_html,
    ".ipynb": load_ipynb,
    ".json": load_json,
}


def load_document(file_path: str) -> list[dict]:
    """Dispatch on file extension. Raises ValueError for unsupported types."""
    ext = Path(file_path).suffix.lower()
    loader = _LOADERS.get(ext)
    if loader is None:
        raise ValueError(f"Unsupported document type: {ext}")
    return loader(file_path)
